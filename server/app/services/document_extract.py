"""
Document text extraction for DLP classification (canonical implementation).

The classification engine scans TEXT. Binary office/PDF attachments decode to
garbage bytes, which is why they previously classified as "Public" and slipped
through — this module turns those bytes into text the classifier can actually
scan (pdf, docx, xlsx, pptx), and decodes the plain-text family directly
(csv, txt, json, xml, ...).

Everything here is defensive on purpose: an attachment is untrusted input.
A parser that blows up, a corrupt/encrypted file, or a format we don't handle
returns a result with ok=False rather than raising, and byte/char caps stop a
hostile attachment from exhausting memory. The caller decides the policy for
un-extractable content — this module never decides allow/block.

Used by /agents/{id}/policy/evaluate when a caller sends `file_content_b64`:
the endpoint agent (USB / file transfer), the browser upload guard, and the
SMTP relay all send raw bytes and let the server do the parsing — so binary
office/PDF attachments are classified on their real text instead of the
compressed bytes (which always looked like "Public" and were let through).
"""
from __future__ import annotations

import io
import logging
import os
from typing import NamedTuple

log = logging.getLogger(__name__)

# Don't even attempt to parse an attachment larger than this.
MAX_EXTRACT_BYTES = 25 * 1024 * 1024

# Cap the text handed to the classifier.
#
# This is a SECURITY boundary, not just a performance knob: text past the cap is
# never scanned, so truncating silently meant "secret at character 1,000,001" =>
# no matches => Public => allowed. Padding a file with filler ahead of the real
# content was therefore a complete bypass (an 87KB zip was enough). Whenever we
# clip we now set Extracted.truncated so the caller can refuse to bless the part
# we never read — see the module docstring's rule: uninspectable != clean.
#
# Sized to stay under ClassificationEngine.MAX_CONTENT_LENGTH (10 MiB) so the
# classifier never silently truncates further. Rule evaluation costs ~700ms per
# 1M chars, putting a worst-case full scan near 7s — inside the endpoint agent's
# default 30s WinHTTP receive timeout.
MAX_TEXT_CHARS = 10_000_000

# Formats that are already text — decode, don't parse.
TEXT_EXTS = {
    ".txt", ".csv", ".tsv", ".log", ".json", ".xml", ".html", ".htm", ".md",
    ".yaml", ".yml", ".ini", ".cfg", ".conf", ".sql", ".eml", ".rtf",
    ".py", ".js", ".ts", ".java", ".c", ".cpp", ".cs", ".go", ".rb", ".php",
    ".sh", ".ps1", ".bat",
}


class Extracted(NamedTuple):
    text: str
    kind: str      # pdf | docx | xlsx | pptx | text | unsupported | error | too_large | empty
    ok: bool       # True when we produced text we trust for classification
    reason: str = ""
    # True when we read only PART of the content (hit MAX_TEXT_CHARS, or an
    # archive hit its safety budget). `text` is still worth classifying — it may
    # convict the file on its own — but the caller must NOT treat a clean result
    # as proof the file is clean, because we did not see all of it.
    truncated: bool = False


def _clip(s: str) -> tuple[str, bool]:
    """Clip to the scan budget. Returns (text, was_truncated)."""
    if len(s) <= MAX_TEXT_CHARS:
        return s, False
    return s[:MAX_TEXT_CHARS], True


def _decode(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return data.decode("utf-8", errors="ignore")


def _looks_textual(text: str) -> bool:
    """True when a decoded blob is mostly printable — i.e. it really is text."""
    if not text:
        return False
    sample = text[:4096]
    printable = sum(c.isprintable() or c.isspace() for c in sample)
    return (printable / max(1, len(sample))) > 0.85


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    if getattr(reader, "is_encrypted", False):
        try:
            reader.decrypt("")  # some PDFs are "encrypted" with an empty owner password
        except Exception:
            raise ValueError("encrypted pdf")
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _extract_docx(data: bytes) -> str:
    import docx
    document = docx.Document(io.BytesIO(data))
    parts = [p.text for p in document.paragraphs]
    for table in document.tables:                      # PII loves living in tables
        for row in table.rows:
            parts.extend(cell.text for cell in row.cells)
    return "\n".join(parts)


def _extract_xlsx(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    try:
        lines = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                if row:
                    lines.append(" ".join("" if v is None else str(v) for v in row))
        return "\n".join(lines)
    finally:
        wb.close()


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                out.append(shape.text_frame.text)
            elif hasattr(shape, "text"):
                out.append(shape.text)
    return "\n".join(out)


_PARSERS = {
    ".pdf": ("pdf", _extract_pdf),
    ".docx": ("docx", _extract_docx),
    ".docm": ("docx", _extract_docx),
    ".xlsx": ("xlsx", _extract_xlsx),
    ".xlsm": ("xlsx", _extract_xlsx),
    ".pptx": ("pptx", _extract_pptx),
}

# Legacy OLE formats need extra tooling we don't ship; call them out explicitly
# so the caller can policy-decide rather than silently treating them as clean.
_LEGACY_EXTS = {".doc", ".xls", ".ppt"}

# ── Archives ────────────────────────────────────────────────────────────────
# Zipping a file is the cheapest possible way to defeat content inspection, so
# we expand archives and classify what's inside. Entries are fed back through
# extract_text(), which means a PDF inside a zip is parsed as a PDF, and a zip
# inside a zip recurses — bounded by the limits below.
#
# Those limits are the zip-bomb defence: a tiny archive can expand to petabytes,
# so we cap nesting depth, entry count and total decompressed bytes, and stop as
# soon as any is hit. Better to inspect the first N MB of a hostile archive than
# to let it exhaust the box.
ARCHIVE_EXTS = {".zip", ".tar", ".tgz", ".gz", ".bz2", ".xz", ".7z"}
MAX_ARCHIVE_DEPTH = 3
MAX_ARCHIVE_ENTRIES = 500
MAX_ARCHIVE_TOTAL_BYTES = 100 * 1024 * 1024


class _ArchiveBudget:
    """Shared across one archive tree so nested members can't each spend the cap."""

    def __init__(self) -> None:
        self.bytes_left = MAX_ARCHIVE_TOTAL_BYTES
        self.entries_left = MAX_ARCHIVE_ENTRIES
        self.truncated = False

    def take(self, n: int) -> bool:
        if self.entries_left <= 0 or n > self.bytes_left:
            self.truncated = True
            return False
        self.entries_left -= 1
        self.bytes_left -= n
        return True


def _extract_zip(data: bytes, depth: int, budget: "_ArchiveBudget") -> tuple[str, bool, str]:
    import zipfile
    parts: list[str] = []
    encrypted = False
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        for info in zf.infolist():
            if info.is_dir():
                continue
            if not budget.take(info.file_size):
                break
            try:
                inner = zf.read(info)
            except RuntimeError as e:              # "password required" => encrypted member
                if "encrypted" in str(e).lower():
                    encrypted = True
                    continue
                raise
            except Exception:
                continue
            sub = extract_text(info.filename, inner, _depth=depth + 1, _budget=budget)
            if sub.truncated:
                budget.truncated = True     # a member we could only partly read
            if sub.text.strip():
                parts.append(sub.text)
    reason = ""
    if encrypted:
        reason = "contains encrypted members"
    return "\n".join(parts), encrypted, reason


def _extract_tar(data: bytes, depth: int, budget: "_ArchiveBudget") -> str:
    import tarfile
    parts: list[str] = []
    with tarfile.open(fileobj=io.BytesIO(data), mode="r:*") as tf:
        for member in tf:
            if not member.isfile():
                continue
            if not budget.take(member.size):
                break
            fh = tf.extractfile(member)
            if fh is None:
                continue
            sub = extract_text(member.name, fh.read(), _depth=depth + 1, _budget=budget)
            if sub.truncated:
                budget.truncated = True
            if sub.text.strip():
                parts.append(sub.text)
    return "\n".join(parts)


def _extract_gz(filename: str, data: bytes, depth: int, budget: "_ArchiveBudget") -> str:
    import gzip
    raw = gzip.decompress(data)
    if not budget.take(len(raw)):
        return ""
    inner_name = filename[:-3] if filename.lower().endswith(".gz") else filename
    sub = extract_text(inner_name, raw, _depth=depth + 1, _budget=budget)
    if sub.truncated:
        budget.truncated = True
    return sub.text


def _extract_7z(data: bytes, depth: int, budget: "_ArchiveBudget") -> str:
    import py7zr                                   # optional dep; absent => unreadable
    parts: list[str] = []
    with py7zr.SevenZipFile(io.BytesIO(data), mode="r") as z:
        for name, bio in (z.readall() or {}).items():
            payload = bio.read()
            if not budget.take(len(payload)):
                break
            sub = extract_text(name, payload, _depth=depth + 1, _budget=budget)
            if sub.truncated:
                budget.truncated = True
            if sub.text.strip():
                parts.append(sub.text)
    return "\n".join(parts)


def sniff_kind(filename: str, data: bytes) -> str:
    """Best-effort format label from the extension, with a magic-byte fallback."""
    ext = os.path.splitext(filename or "")[1].lower()
    if ext in _PARSERS:
        return _PARSERS[ext][0]
    if ext in TEXT_EXTS:
        return "text"
    if ext in _LEGACY_EXTS:
        return "unsupported"
    if data[:5] == b"%PDF-":
        return "pdf"
    if data[:4] == b"PK\x03\x04":
        return "ooxml"   # a zip container — needs the extension to disambiguate
    return "unknown"


def extract_text(filename: str, data: bytes, _depth: int = 0,
                 _budget: "_ArchiveBudget | None" = None) -> Extracted:
    """Return text suitable for classification, plus how we got it.

    _depth/_budget are internal: they bound recursion and total decompression
    when expanding archives (see ARCHIVE_EXTS).
    """
    if not data:
        return Extracted("", "empty", True, "no content")
    if len(data) > MAX_EXTRACT_BYTES:
        return Extracted("", "too_large", False, f"{len(data)} bytes exceeds cap")

    ext = os.path.splitext(filename or "")[1].lower()

    # Office parsers FIRST: .docx/.xlsx/.pptx are themselves zip containers, so
    # they must go to their own parser rather than the archive expander (which
    # would walk their internal XML plumbing instead of reading the document).
    parser = _PARSERS.get(ext)
    if parser is None and data[:5] == b"%PDF-":
        parser = _PARSERS[".pdf"]                       # mislabeled/extension-less PDF

    if parser is None:
        is_archive_ext = ext in ARCHIVE_EXTS
        looks_zip = data[:4] == b"PK\x03\x04"
        if is_archive_ext or looks_zip:
            if _depth >= MAX_ARCHIVE_DEPTH:
                return Extracted("", "archive", False, "nesting depth limit reached")
            budget = _budget or _ArchiveBudget()
            try:
                if ext in (".tar", ".tgz", ".bz2", ".xz") or (ext == ".gz" and filename.lower().endswith(".tar.gz")):
                    text = _extract_tar(data, _depth, budget)
                    encrypted, reason = False, ""
                elif ext == ".gz":
                    text = _extract_gz(filename, data, _depth, budget)
                    encrypted, reason = False, ""
                elif ext == ".7z":
                    text = _extract_7z(data, _depth, budget)
                    encrypted, reason = False, ""
                else:
                    text, encrypted, reason = _extract_zip(data, _depth, budget)
            except ImportError:
                return Extracted("", "archive", False, f"{ext} support not installed")
            except Exception as e:                      # corrupt / encrypted / unsupported
                # Password-protected archives raise deep library-specific errors;
                # surface a clean, actionable reason instead of a raw dump.
                msg = str(e).lower()
                if "password" in msg or "encrypted" in msg:
                    log.info("encrypted archive %s — cannot inspect", filename)
                    return Extracted("", "archive", False, "encrypted archive (password required)")
                log.warning("archive extract failed for %s: %s", filename, e)
                return Extracted("", "archive", False, f"archive: {str(e)[:120]}")

            if budget.truncated:
                reason = (reason + "; " if reason else "") + "archive truncated at safety limits"
            if text.strip():
                # Found readable content — classify it. Note ok=True even if some
                # members were encrypted: what we COULD read still counts.
                clipped, was_clipped = _clip(text)
                if was_clipped:
                    reason = (reason + "; " if reason else "") + \
                        f"archive text clipped at {MAX_TEXT_CHARS} chars"
                # budget.truncated => we stopped early (entry/byte/depth limit or
                # a partly-read member), so the archive was NOT fully inspected.
                return Extracted(clipped, "archive", True, reason,
                                 was_clipped or budget.truncated)
            # Nothing readable: an encrypted or opaque archive. Say so rather
            # than letting empty text classify as "Public".
            return Extracted("", "archive", False, reason or "no readable content in archive")

    if parser is not None:
        kind, fn = parser
        try:
            text, was_clipped = _clip(fn(data))
        except Exception as e:                          # corrupt / encrypted / extension lies
            # SECURITY: never conclude "unreadable => clean" just because the
            # parser for the *claimed* format rejected it. The commonest cause
            # is that the extension is a lie — rename secret.txt to secret.docx
            # and the Word parser fails; if we stopped here the content would
            # classify as Public and sail through. So fall back to scanning the
            # raw bytes as text, which catches exactly that evasion.
            fallback = _decode(data)
            if _looks_textual(fallback):
                log.warning("%s: %s parser failed (%s) — content is textual, scanning as text",
                            filename, kind, e)
                clipped, clipped_flag = _clip(fallback)
                return Extracted(clipped, "text", True,
                                 f"{kind} parse failed; content scanned as text", clipped_flag)
            log.warning("extract failed for %s (%s): %s", filename, kind, e)
            return Extracted("", "error", False, f"{kind}: {e}")
        if not text.strip():
            # Parsed fine but produced nothing — e.g. a scanned/image-only PDF.
            return Extracted("", kind, False, "no extractable text (image-only?)")
        reason = f"text clipped at {MAX_TEXT_CHARS} chars" if was_clipped else ""
        return Extracted(text, kind, True, reason, was_clipped)

    if ext in TEXT_EXTS:
        clipped, was_clipped = _clip(_decode(data))
        return Extracted(clipped, "text", True,
                         f"text clipped at {MAX_TEXT_CHARS} chars" if was_clipped else "",
                         was_clipped)

    if ext in _LEGACY_EXTS:
        return Extracted("", "unsupported", False, f"legacy format {ext} not supported")

    # Unknown extension: if it decodes cleanly it's text; otherwise say so.
    text = _decode(data)
    if _looks_textual(text):
        clipped, was_clipped = _clip(text)
        return Extracted(clipped, "text", True,
                         f"text clipped at {MAX_TEXT_CHARS} chars" if was_clipped else "",
                         was_clipped)

    return Extracted("", "unsupported", False, f"binary/unknown format {ext or '(none)'}")
